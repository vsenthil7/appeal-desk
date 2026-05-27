"""
Build the Appeal-Desk pitch deck for the Reddit Developer Platform
hackathon (AT-Hack0022) submission.

12 slides, 16:9 widescreen, Appeal-Desk brand:
  - Navy #0c1a2b background (authority, due-process)
  - Cream #f5efe3 primary text (humane, civil tone)
  - Slate #6b7a8f secondary text
  - Verdigris #4caf7d accent (overturned / approved)
  - Coral #e26464 accent (upheld / refused)

Output (relative to appeal-desk/ repo root):
  - submission_media/appeal-desk-pitch-deck-{stamp}.pptx
  - submission_media/appeal-desk-pitch-deck-{stamp}.pdf  (produced separately by pptx_to_pdf.ps1)
  - submission_media/_backup/{same files}

Content is grounded in the actual repo:
  - 288 tests pass / 18 test files
  - Coverage: stmts 99.97 / branch 99.06 / funcs 100 / lines 99.97 (gates pass)
  - 4 build commits on main (vsenthil7/appeal-desk)
  - AI is *assistive only*; product works fully with AI stripped out
  - Numbers verified by AppealDesk-Code-Review.md + CODE-FIX-NOTES.md

Ported from the ATRIO submission pipeline (atrio/scripts/build_pitch_deck.py).
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ---------- Brand ----------
NAVY = RGBColor(0x0C, 0x1A, 0x2B)            # background
CREAM = RGBColor(0xF5, 0xEF, 0xE3)           # primary text on dark
TEXT_PRIMARY = RGBColor(0xE5, 0xDF, 0xD3)    # body text on dark
TEXT_SECONDARY = RGBColor(0x9B, 0xA8, 0xBA)  # muted
VERDIGRIS = RGBColor(0x4C, 0xAF, 0x7D)       # accent: overturned/approved
CORAL = RGBColor(0xE2, 0x64, 0x64)           # accent: upheld/refused
BLUE = RGBColor(0x5E, 0x8B, 0xD4)            # informational accent
PANEL_BG = RGBColor(0x14, 0x26, 0x3E)        # slightly lifted panel against NAVY


# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HACKATHON_LABEL = "AT-HACK0022  \u00b7  REDDIT DEV PLATFORM 2026"
REPO_FOOTER = "Appeal-Desk  \u00b7  github.com/vsenthil7/appeal-desk  \u00b7  MIT"
TOTAL_SLIDES = 12


def fill_slide(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    italic=False,
    color=CREAM,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Inter",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_top_strip(slide, hackathon=HACKATHON_LABEL, color=VERDIGRIS):
    add_rect(slide, Inches(0.5), Inches(0.35), Inches(0.25), Inches(0.06), color)
    add_text(
        slide,
        Inches(0.85), Inches(0.25), Inches(10), Inches(0.3),
        hackathon,
        size=10, color=TEXT_SECONDARY, font="Inter",
    )


def add_footer(slide, page_num=None, total=TOTAL_SLIDES):
    add_text(
        slide,
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
        REPO_FOOTER,
        size=9, color=TEXT_SECONDARY, font="Inter",
    )
    if page_num is not None:
        add_text(
            slide,
            Inches(11), Inches(7.05), Inches(2), Inches(0.3),
            f"{page_num} / {total}",
            size=9, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.RIGHT,
        )


# ---------- Slide builders ----------


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(1.5),
             "APPEAL-DESK", size=84, bold=True, color=CREAM, font="Inter")
    add_rect(s, Inches(0.55), Inches(3.55), Inches(1.5), Inches(0.06), VERDIGRIS)
    add_text(s, Inches(0.5), Inches(3.85), Inches(12), Inches(1.0),
             "A fair appeals desk for modteams.",
             size=36, italic=True, color=CREAM, font="Inter")
    add_text(s, Inches(0.5), Inches(4.85), Inches(12), Inches(1.0),
             "Structured intake \u00b7 audit-trail \u00b7 one-tap decisions \u00b7 civil replies.",
             size=22, color=TEXT_PRIMARY, font="Inter")
    add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
             "Reddit Developer Platform Hackathon 2026  \u00b7  AT-Hack0022  \u00b7  github.com/vsenthil7/appeal-desk",
             size=12, color=TEXT_SECONDARY, font="Inter")
    return s


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The problem", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Appeals live in modmail. Justice doesn't scale that way.",
             size=36, bold=True, color=CREAM, font="Inter")
    add_text(s, Inches(0.5), Inches(2.8), Inches(12), Inches(2.5),
             "Today, a banned or removed user has one option: write to modmail and hope.\n\n"
             "  \u2192 Mods get an unstructured wall of text \u2014 no link back to the action.\n"
             "  \u2192 No history. The same user can re-appeal the same ban 6 times.\n"
             "  \u2192 No audit trail. \"Why did we uphold this?\" is unanswerable a month later.\n"
             "  \u2192 No civility scaffold. Replies are drafted from scratch under pressure.\n\n"
             "The result: inconsistent decisions, mod burnout, and users who feel unheard\n"
             "even when their appeal is fair.",
             size=16, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 2)
    return s


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The solution", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "A dedicated, fair appeals desk \u2014 built into the subreddit.",
             size=30, bold=True, color=CREAM, font="Inter")

    cols = [
        ("STRUCTURE",
         "Ban or removal happens \u2192 the user sees a structured\nintake form linked to the exact action. No more\nfreeform modmail. No more lost context."),
        ("ASSIST",
         "Mods get a queue with full history, deterministic\nduplicate detection, and templated civil replies.\nThree buttons: uphold, overturn, ask for more."),
        ("AUDIT",
         "Every decision, reply, and state change is logged.\nGDPR-style erasure preserves a tombstone.\nRetention is enforced, not just promised."),
    ]
    col_w = Inches(4.1)
    gap = Inches(0.1)
    for i, (header, body) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        add_rect(s, x, Inches(2.8), col_w, Inches(0.06), VERDIGRIS)
        add_text(s, x, Inches(2.95), col_w, Inches(0.4),
                 header, size=14, bold=True, color=VERDIGRIS, font="Inter")
        add_text(s, x, Inches(3.45), col_w, Inches(3),
                 body, size=14, color=TEXT_PRIMARY, font="Inter")

    add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
             "AI is assistive only \u2014 a tone-softener and triage hint that the mod can ignore.\n"
             "The product works end-to-end with AI stripped out. This is enforced by `selectProvider()` and tested.",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 3)
    return s


def slide_user_journey(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The user journey", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "From mod action to civil reply \u2014 in five tracked steps.",
             size=28, bold=True, color=CREAM, font="Inter")

    steps = [
        ("1", "ACTION",   "Mod bans or removes content. Appeal-Desk snapshots the original\nremoval reason at the moment of action \u2014 user can't edit it later."),
        ("2", "INVITE",   "User gets a civil modmail with a link to the appeal form. Form is\npre-filled with the action context (read-only) plus a reason field."),
        ("3", "INTAKE",   "Submission is validated, sanitised, rate-limited, and deduplicated\nagainst the user's prior appeals. AI triage is optional, never blocking."),
        ("4", "REVIEW",   "Mods see a dashboard with full history, near-duplicate flags, and\nthree one-tap buttons. Each button opens a reply-confirm form."),
        ("5", "DECIDE",   "Mod approves the reply. Decision is recorded *before* the reply\nis sent \u2014 transient send failures never roll back the decision."),
    ]
    for i, (num, title, body) in enumerate(steps):
        y = Inches(2.7) + Inches(0.85) * i
        add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.85),
                 num, size=36, bold=True, color=VERDIGRIS, font="Inter")
        add_text(s, Inches(1.3), y + Inches(0.05), Inches(11), Inches(0.4),
                 title, size=15, bold=True, color=CREAM, font="Inter")
        add_text(s, Inches(1.3), y + Inches(0.4), Inches(11.4), Inches(0.5),
                 body, size=11, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 4)
    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Architecture", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Platform-free core. Devvit-thin shell.",
             size=30, bold=True, color=CREAM, font="Inter")

    tiers = [
        ("CORE",  "TypeScript \u00b7 zero Devvit imports",
                  "store \u00b7 service \u00b7 concurrency \u00b7 validation \u00b7 dedup \u00b7 templates \u00b7 retention",
                  VERDIGRIS),
        ("AI",    "Optional provider behind `selectProvider()`",
                  "NoopAiProvider is the default \u2014 product is fully functional with AI off",
                  BLUE),
        ("SHELL", "Devvit (.tsx) wiring only",
                  "triggers \u00b7 intake form \u00b7 dashboard custom post \u00b7 menu \u00b7 scheduler",
                  CORAL),
    ]
    for i, (label, stack, note, color) in enumerate(tiers):
        y = Inches(2.8) + Inches(1.2) * i
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), PANEL_BG, line_color=color)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 label, size=16, bold=True, color=color, font="Inter")
        add_text(s, Inches(3.5), y + Inches(0.1), Inches(9), Inches(0.4),
                 stack, size=14, bold=True, color=CREAM, font="Inter")
        add_text(s, Inches(3.5), y + Inches(0.55), Inches(9), Inches(0.4),
                 note, size=11, color=TEXT_SECONDARY, font="Inter", italic=True)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "The platform-free core has 100% test coverage *without* the Devvit runtime. "
             "That's the headline architectural call.",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 5)
    return s


def slide_concurrency(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Concurrency \u00b7 The hard part, done", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "One ban. One appeal. No matter how many tabs are open.",
             size=28, bold=True, color=CREAM, font="Inter")

    states = [
        ("submitted",   VERDIGRIS),
        ("in_review",   BLUE),
        ("awaiting_user", BLUE),
        ("resolved",    CORAL),
    ]
    box_w = Inches(2.6)
    box_h = Inches(0.7)
    gap = Inches(0.4)
    x0 = Inches(0.5)
    y0 = Inches(2.9)
    for i, (state, color) in enumerate(states):
        x = x0 + (box_w + gap) * i
        add_rect(s, x, y0, box_w, box_h, PANEL_BG, line_color=color)
        add_text(s, x, y0 + Inches(0.18), box_w, Inches(0.4),
                 state, size=13, bold=True, color=CREAM, font="Inter", align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, x + box_w - Inches(0.05), y0 + Inches(0.12), gap + Inches(0.2), Inches(0.5),
                     "\u2192", size=24, bold=True, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.CENTER)

    guarantees = [
        ("\u2713 Atomic action-lock",
         "WATCH/MULTI/EXEC CAS. Two users opening the same appeal in the same\nmillisecond \u2014 only one wins. Proven by parallel-execution tests."),
        ("\u2713 Terminal-state guard",
         "A `resolved` appeal rejects further decisions with INVALID_STATE_TRANSITION.\nNo accidental over-write of a final verdict."),
        ("\u2713 Record-then-send",
         "Decision is persisted *before* the reply is dispatched. A modmail send\nfailure throws REPLY_DELIVERY_FAILED but the decision stands."),
        ("\u2713 Tie-safe pagination",
         "Cursor is `{score, id}`, not a bare score. Same-millisecond appeals\nno longer skip pages \u2014 fix landed in BUILD 003."),
    ]
    for i, (gate, detail) in enumerate(guarantees):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + Inches(6.2) * col
        y = Inches(4.2) + Inches(1.0) * row
        add_text(s, x, y, Inches(6.0), Inches(0.4),
                 gate, size=13, bold=True, color=VERDIGRIS, font="Inter")
        add_text(s, x, y + Inches(0.4), Inches(6.0), Inches(0.6),
                 detail, size=11, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 6)
    return s


def slide_ai_stance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The AI stance", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "AI drafts. Humans decide. The product survives AI removal.",
             size=28, bold=True, color=CREAM, font="Inter")

    add_text(s, Inches(0.5), Inches(2.6), Inches(12), Inches(0.5),
             "Three rules, enforced structurally:",
             size=14, color=CREAM, bold=True, font="Inter")

    rules = [
        ("1 \u00b7 AI never blocks intake",
         "`submitAppeal()` calls the triage provider best-effort. Failure returns null.\nThe appeal is created either way. Tested with a failing provider."),
        ("2 \u00b7 The mod always confirms",
         "Every one-tap decision opens a reply-confirm form. The drafted reply is\neditable. AI output is clamped (empty / >4\u00d7 length \u2192 fallback to template)."),
        ("3 \u00b7 NoopAiProvider is the default",
         "`selectProvider(aiEnabled, backend)` returns a no-op when AI is off or no\nbackend is wired. The full happy path works without an API key."),
    ]
    for i, (rule, detail) in enumerate(rules):
        y = Inches(3.2) + Inches(1.1) * i
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), PANEL_BG, line_color=VERDIGRIS)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(11.5), Inches(0.4),
                 rule, size=14, bold=True, color=CREAM, font="Inter")
        add_text(s, Inches(0.7), y + Inches(0.5), Inches(11.5), Inches(0.5),
                 detail, size=11, color=TEXT_PRIMARY, font="Inter")

    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "Why this matters: a tool that *needs* AI fails closed if the model has a bad day. "
             "Appeal-Desk fails open \u2014 a slightly less polished reply, but the workflow still works.",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 7)
    return s


def slide_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Proof, not promises", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Numbers from the actual repo. Reproducible on your machine.",
             size=28, bold=True, color=CREAM, font="Inter")

    metrics = [
        ("288 / 288", "tests pass"),
        ("99.97 %",   "line coverage"),
        ("100 %",     "function coverage"),
        ("99.06 %",   "branch coverage"),
        ("18",        "test files"),
        ("0",         "TypeScript errors"),
        ("0",         "ESLint issues"),
        ("4",         "build commits on main"),
    ]
    box_w = Inches(3.0)
    box_h = Inches(1.4)
    for i, (big, small) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + (box_w + Inches(0.1)) * col
        y = Inches(2.9) + (box_h + Inches(0.2)) * row
        add_rect(s, x, y, box_w, box_h, PANEL_BG, line_color=VERDIGRIS)
        add_text(s, x, y + Inches(0.2), box_w, Inches(0.6),
                 big, size=28, bold=True, color=CREAM, font="Inter", align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.85), box_w, Inches(0.45),
                 small, size=11, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
             "`npm ci && npm run lint && npx tsc --noEmit && npx vitest run --coverage`",
             size=11, color=TEXT_SECONDARY, font="Consolas", italic=True)
    add_footer(s, 8)
    return s


def slide_findings(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Honest defect log", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Reviewed end-to-end. Each defect found \u2014 and fixed.",
             size=28, bold=True, color=CREAM, font="Inter")

    findings = [
        ("F1 \u00b7 HIGH",   "`npm run lint` broken", "No ESLint config existed.",
         "Added `.eslintrc.cjs` + deps. Clean."),
        ("F2 \u00b7 MED",    "Same-ms cursor tie-skip", "`cursor - 1` dropped co-scored entries.",
         "Tuple cursor `{score, id}`. New test."),
        ("F3 \u00b7 MED",    "Unbounded queue read",   "`openQueuePage` hydrated whole index.",
         "`limit: {offset, count}` everywhere."),
        ("F4 \u00b7 LOW",    "Seed key smuggled",       "`actionLock` builder reused.",
         "First-class `keys.actionSeed()`."),
        ("F5 \u00b7 LOW",    "purgeExpired over-fetch", "Symmetric with F3.",
         "Same bounded-read treatment."),
        ("F6 \u00b7 NIT",    "Stray `ms` field",        "Type cast hid a tiny lie.",
         "Dropped the field and the cast."),
        ("Wiring",   "syncConfig only on install", "Settings edits ignored till reinstall.",
         "Now runs on submission + AppUpgrade."),
        ("Wiring",   "Retention was dead code",  "purgeExpired/redactAppeal never called.",
         "Daily scheduler job + service entry points."),
    ]
    row_h = Inches(0.5)
    for i, (severity, title, problem, fix) in enumerate(findings):
        y = Inches(2.7) + row_h * i
        add_text(s, Inches(0.5),  y, Inches(1.3),  Inches(0.4),
                 severity, size=10, bold=True, color=CORAL, font="Inter")
        add_text(s, Inches(2.0),  y, Inches(3.0),  Inches(0.4),
                 title,    size=11, bold=True, color=CREAM, font="Inter")
        add_text(s, Inches(5.1),  y, Inches(3.6),  Inches(0.4),
                 problem,  size=10, color=TEXT_SECONDARY, font="Inter", italic=True)
        add_text(s, Inches(8.8),  y, Inches(4.0),  Inches(0.4),
                 fix,      size=10, color=VERDIGRIS, font="Inter")
    add_footer(s, 9)
    return s


def slide_what_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Beyond the MVP", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "What ships next \u2014 width first, then depth.",
             size=28, bold=True, color=CREAM, font="Inter")

    columns = [
        ("Width (high-leverage modules)", [
            "core/analytics \u2014 overturn rate, time-to-decision",
            "core/policy \u2014 eligibility + cooldowns per sub",
            "core/notifications \u2014 user-side status updates",
            "core/escalation \u2014 multi-mod sign-off",
            "core/export \u2014 CSV/JSON portability",
            "core/i18n \u2014 message catalogs per sub language",
        ]),
        ("Depth (richer existing modules)", [
            "Bulk decisions on near-duplicates",
            "Incremental dedup (MinHash, no full re-hydration)",
            "Conditional templates + token linter",
            "PII pre-flags + Unicode normalisation in dedup",
            "Correlation IDs in observability",
            "AI eval harness with golden fixtures",
        ]),
    ]
    for i, (header, items) in enumerate(columns):
        x = Inches(0.5) + Inches(6.4) * i
        add_rect(s, x, Inches(2.7), Inches(6.0), Inches(0.05), VERDIGRIS)
        add_text(s, x, Inches(2.85), Inches(6.0), Inches(0.4),
                 header, size=15, bold=True, color=CREAM, font="Inter")
        for j, item in enumerate(items):
            add_text(s, x, Inches(3.35) + Inches(0.45) * j, Inches(6.0), Inches(0.4),
                     "\u00b7  " + item, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 10)
    return s


def slide_install(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Install \u00b7 Try it", size=12, color=VERDIGRIS, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Two commands. Live on your test sub in 90 seconds.",
             size=28, bold=True, color=CREAM, font="Inter")

    add_rect(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(2.2), PANEL_BG, line_color=VERDIGRIS)
    add_text(s, Inches(0.7), Inches(2.9), Inches(12), Inches(0.4),
             "From source:", size=13, bold=True, color=VERDIGRIS, font="Inter")
    add_text(s, Inches(0.7), Inches(3.3), Inches(12), Inches(0.4),
             "$ git clone https://github.com/vsenthil7/appeal-desk && cd appeal-desk",
             size=13, color=CREAM, font="Consolas")
    add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4),
             "$ npm ci && devvit upload && devvit install r/your_test_sub",
             size=13, color=CREAM, font="Consolas")
    add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.4),
             "From the App Directory:", size=13, bold=True, color=VERDIGRIS, font="Inter")
    add_text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
             "developers.reddit.com/apps/appeal-desk  \u2192  Install on subreddit",
             size=13, color=CREAM, font="Consolas")

    add_text(s, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
             "Permissions requested:", size=14, bold=True, color=CREAM, font="Inter")
    perms = [
        "reddit \u2014 read mod actions, send modmail (mod scope)",
        "redis \u2014 store appeals, history, locks, indexes",
        "scheduler \u2014 SLA nudges + daily retention purge",
        "http: false \u2014 fully on-platform, no external services",
    ]
    for i, p in enumerate(perms):
        add_text(s, Inches(0.7), Inches(5.85) + Inches(0.32) * i, Inches(12), Inches(0.4),
                 "\u00b7  " + p, size=11, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 11)
    return s


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, NAVY)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(1.4),
             "Appeal-Desk",
             size=64, bold=True, color=CREAM, font="Inter")
    add_rect(s, Inches(0.55), Inches(3.4), Inches(1.5), Inches(0.06), VERDIGRIS)
    add_text(s, Inches(0.5), Inches(3.65), Inches(12), Inches(1.0),
             "A fair appeals desk for modteams.",
             size=28, italic=True, color=CREAM, font="Inter")
    add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(1.0),
             "AI assists. Humans decide. The audit trail proves it.",
             size=18, color=TEXT_PRIMARY, font="Inter")

    add_text(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
             "github.com/vsenthil7/appeal-desk",
             size=14, color=VERDIGRIS, font="Consolas")
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
             "MIT  \u00b7  288/288 tests  \u00b7  99.97 % coverage  \u00b7  4 BUILD commits  \u00b7  Devvit App Directory",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "Thank you.",
             size=12, color=TEXT_SECONDARY, font="Inter")
    return s


# ---------- Main ----------


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,         # 1
        slide_problem,       # 2
        slide_solution,      # 3
        slide_user_journey,  # 4
        slide_architecture,  # 5
        slide_concurrency,   # 6
        slide_ai_stance,     # 7
        slide_proof,         # 8
        slide_findings,      # 9
        slide_what_next,     # 10
        slide_install,       # 11
        slide_closing,       # 12
    ]
    for b in builders:
        b(prs)

    out_dir = Path(__file__).resolve().parent.parent / "submission_media"
    out_dir.mkdir(parents=True, exist_ok=True)
    backup_dir = out_dir / "_backup"
    backup_dir.mkdir(exist_ok=True)

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_pptx = out_dir / f"appeal-desk-pitch-deck-{stamp}.pptx"
    prs.save(out_pptx)

    import shutil
    shutil.copy2(out_pptx, backup_dir / out_pptx.name)

    print(f"deck: {out_pptx} ({out_pptx.stat().st_size / 1024:.1f} KB)")
    return out_pptx


if __name__ == "__main__":
    build_deck()
